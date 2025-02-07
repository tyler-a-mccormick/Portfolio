# -*- coding: utf-8 -*-
"""
Created on Fri Apr 19 13:49:08 2024

@author: McCormickT
"""

import pandas as pd
import numpy as np
from sklearn.model_selection import train_test_split
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import classification_report, confusion_matrix, roc_auc_score
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF

# Load the dataset
data = pd.read_csv("I:\spy_mods\Extract 2.csv")

# Fill missing values
data.fillna("Missing", inplace=True)

# Feature engineering: combining features and binning
data["utm_campaign_category"] = data["ssl_utm_campaign"].astype(str) + "_" + data["ssl_typeoffeaturedname"].astype(str)
data["decisiontimefrom_binned"] = pd.qcut(data["Sum of ssl_decisiontimefrom"], 4, labels=False, duplicates="drop")

# Define data for modeling
x = data[['utm_campaign_category', 'decisiontimefrom_binned', 'ssl_utm_source', 'ssl_utm_medium', 'ssl_utm_term']]
y = data['ssl_firstTrue']

# Train a Random Forest Classifier
x_train, x_test, y_train, y_test = train_test_split(x, y, test_size=0.3, random_state=42)

forest = RandomForestClassifier(n_estimators=100, random_state=42)
forest.fit(x_train, y_train)

# Predict and evaluate the model
y_pred = forest.predict(x_test)
print(classification_report(y_test, y_pred))
print("ROC-AUC Score:", roc_auc_score(y_test, forest.predict_proba(x_test)[:, 1]))

# Feature importance visualization
importances = forest.feature_importances_
indices = np.argsort(importances)[::-1]

plt.figure(figsize=(10, 6))
plt.title("Feature Importances")
plt.bar(range(x_train.shape[1]), importances[indices], color="skyblue", align="center")
plt.xticks(range(x_train.shape[1]), x_train.columns[indices], rotation=90)
plt.tight_layout()
plt.show()

# Export findings to PowerPoint
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Feature Importance Analysis"

# Add the plot to the slide
img_path = "feature_importance.png"
plt.savefig(img_path)
left = Inches(1)
top = Inches(1)
slide.shapes.add_picture(img_path, left, top, width=Inches(8), height=Inches(4.5))

prs.save("Feature_Importance_Presentation.pptx")

# Export findings to PDF
pdf = FPDF()
pdf.add_page()
pdf.set_font("Arial", size=12)
pdf.cell(200, 10, txt="Feature Importance Analysis", ln=1, align='C')
pdf.output("Analysis_Report.pdf")
